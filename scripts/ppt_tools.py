#!/bin/env/python
# -*- coding:utf-8 -*-

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt



# _ppt = Presentation()

# ===================================
def add_4pngs_to_1page(ppt, pngs, title, labels=[]):
    ''' '''
    _ppt=ppt

    width = _ppt.slide_width
    height = _ppt.slide_height
    # print " width="+str(width)+ "  height="+str(height)

    picsize = int((height-Pt(32))/2) if len(pngs)>2 else int((width/2))
    if len(pngs) > 2 : 
        hoffset = int((width/2 - picsize)/2)
        voffset = int(height/2-picsize)

        lefttop=[ [hoffset, 2*voffset], [width/2+hoffset, 2*voffset],
              [hoffset, height/2+voffset], [width/2+hoffset, height/2+voffset] ]

    else:
        voffset = int( (height-picsize)/2 )
        lefttop = [ [0, voffset], [picsize, voffset] ]

    blank_slide_layout = _ppt.slide_layouts[6]
    isub = 0
    slide = _ppt.slides.add_slide(blank_slide_layout)
    add_title_line(slide, title, width/10, 0, Pt(24))

    for ip in range(0, len(pngs)):
        if isub > 3:
            slide = _ppt.slides.add_slide(blank_slide_layout)
            isub = 0
        pic = slide.shapes.add_picture(pngs[ip], 0, 0)

        pic.width = picsize
        pic.height = picsize
        pic.left = lefttop[isub][0]
        pic.top = lefttop[isub][1]
        print "Adding "+pngs[ip]+str(pic.width)+"  " + str(pic.height)

        if len(labels) >= 0:
            add_title_line(slide, labels[ip], pic.left+picsize/2, pic.top+3*Pt(14), Pt(14))

        isub += 1


# ===================================
def add_pngpage(ppt, funit): 
    '''  funit=71(81), 72, 73(82), 74, 75(83), 76  '''

    # global _ppt
    _ppt=ppt

    width = _ppt.slide_width
    height = _ppt.slide_height
    # print " width="+str(width)+ "  height="+str(height)

    picsize = int(0.49*float(height))
    hoffset = int((width/2 - picsize)/2)
    voffset = int(height/2-picsize)

    lefttop=[ [hoffset, voffset], [width/2+hoffset, voffset], 
              [hoffset, height/2+voffset], [width/2+hoffset, height/2+voffset] ]

    period = ["1s", "1M", "1h", "1d", "1w", "1m", "3m", "1y", "4y", "Xy", "Zy"]
    pname = ["1 second", "1 minute", "1 hour", "1 day", "1 week", "1 month", "3 months", 
                   "1 year", "4 years", "10 years", "50 years"]
    figkeys={"71":[81, "doseeq", "All"],  "72":[0, "activity","All"], 
             "73":[82, "doseeq", "Allm"], "74":[0, "activity","Allm"], 
             "75":[83, "doseeq", "Allt"], "76":[0, "activity","Allt"],
             "81":[0, "doseeq", "All"], "82":[0, "doseeq","mid"], "83":[0, "doseeq","tar"]}
    figdir="figs-trim"

    # Blank slide layout
    blank_slide_layout = _ppt.slide_layouts[6]
    isub=0

    slide = _ppt.slides.add_slide(blank_slide_layout)  
    
    su=str(funit)
    if figkeys[su][0] != 0: 
        psu = str(figkeys[su][0])
        fname="%s/f%s-%s-pri-%s.png" % (figdir, psu, figkeys[psu][1], figkeys[psu][2])
        pic = slide.shapes.add_picture(fname, 0, 0)
        pic.width = picsize
        pic.height = picsize
        pic.left = lefttop[0][0]
        pic.top = lefttop[0][1]
        isub += 1

    for ip in range(0, len(period)):
        fname="%s/f%s-%s-%s-%s.png" % (figdir, su, figkeys[su][1], period[ip], figkeys[su][2])
        if isub > 3: 
            slide = _ppt.slides.add_slide(blank_slide_layout)  
            isub = 0
        pic = slide.shapes.add_picture(fname, 0, 0)


        pic.width = picsize
        pic.height = picsize
        pic.left = lefttop[isub][0]
        pic.top = lefttop[isub][1]
        print "Adding "+fname+str(pic.width)+"  " + str(pic.height)
        isub += 1


# ===================================
def add_title_line(slide, message, left, top, font_size):

    # global _ppt

    txBox=slide.shapes.add_textbox(left, top, len(message)*font_size, font_size)

    txBox.text_frame.add_paragraph()
    tf = txBox.text_frame.paragraphs[0]
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = False
    tf.margin_top = 0
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    run = tf.add_run()
    run.text = message
    font = run.font
    font.name = 'Calibri'
    font.size = font_size
    font.bold = None
    font.italic = None  # cause value to be inherited from theme
    # font.color.theme_color = MSO_THEME_COLOR.ACCENT_5


#  ===================================
def add_textmessage(ppt, messages):

    # global _ppt
    _ppt = ppt

    blank_slide_layout = _ppt.slide_layouts[6]
    slide = _ppt.slides.add_slide(blank_slide_layout)
    width = _ppt.slide_width
    height = _ppt.slide_height

    txBox=slide.shapes.add_textbox(width/8, height/4, 50*Pt(24), Pt(48))

    for it in range(0, len(messages)):
        txBox.text_frame.add_paragraph()
        tf = txBox.text_frame.paragraphs[it]
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = False
        tf.margin_top = 0
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        run = tf.add_run()
        run.text = messages[it]
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(32)
        font.bold = None
        font.italic = None  # cause value to be inherited from theme
        # font.color.theme_color = MSO_THEME_COLOR.ACCENT_5

        
# ===================================
if __name__ == "__main__":

    ppt=Presentation()

    # Comparison amoung running 
    # Adding text page        
    add_textmessage(ppt, ["Comparison of Dose-EQ", 
            "Beam run 1 year, 4 years, 20 years", 
            "2625 Bxs, 5Hz, 5000 hours/year" ])
    pdata =[ ["dose_vs_r", "DoseEq vs R"] ]
    ftype={"71":"All", "73":"Allm", "75":"Allt"}
    for iu in [71, 73, 75]:
        pdata += [ ["f%d-doseeq-1s-%s"%(iu, ftype[str(iu)]), "Cooling 1 second"],
                 ["f%d-doseeq-1d-%s"%(iu, ftype[str(iu)]), "Cooling 1 day"],
                 ["f%d-doseeq-1m-%s"%(iu, ftype[str(iu)]), "Cooling 1 month"],
                 ["f%d-doseeq-1y-%s"%(iu, ftype[str(iu)]), "Cooling 1 year"],
                 ["f%d-doseeq-Xy-%s"%(iu, ftype[str(iu)]), "Cooling 10 years"]]

    for pd in pdata:
        pngs_dose_vs_r=[]
        for tdir in ["1year", "4year", "20year"]:
           pngs_dose_vs_r.append(tdir+"/figs-trim/%s.png" % pd[0])
        add_4pngs_to_1page(ppt, pngs_dose_vs_r, "%s, beam run 1, 4, 20 years" % pd[1], 
            labels=["1 year run", "4 years run ", "20 years run "])

    

    for tdir in ["1year", "4year", "20year"]:
    
        os.chdir(tdir)

        # Adding text page        
        add_textmessage(ppt, [tdir+" Dose-EQ", 
            "All region, Upto middle region, Target region"])

        for iu in [71, 73, 75]:
            add_pngpage(ppt, iu)

        
        add_textmessage(ppt, [tdir+" Activity", 
            "All region, Upto middle region, Target region"])

        for iu in [72, 74, 76]:
            add_pngpage(ppt, iu)


        os.chdir("..")

    _ppt.save("figure.pptx")

